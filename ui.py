# app.py

import streamlit as st

st.set_page_config(
    page_title="💎 Rishab Gems ‒ Diamond Jewellery Invoice Generator",
    layout="wide",
)

import pandas as pd
from io import BytesIO
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, MSO_VERTICAL_ANCHOR

# ──────────────────────────────────────────────────────────────────────────────
#  Helper: Fill PPTX template with invoice data (finding tables by alt_text)
# ──────────────────────────────────────────────────────────────────────────────
def generate_filled_invoice(rows, template_path):
    """
    1. Open the PPTX template.
    2. Locate the table whose name == "LineItems", write each data row BELOW the header.
    3. Locate the table whose name == "BillingSummary", fill Subtotal/Rounding/NET PAYABLE.
    4. Center-align every cell’s text and preserve Poppins 12 pt.
    5. Return the PPTX bytes.
    """
    prs = Presentation(template_path)
    slide = prs.slides[0]

    # 1) Find both tables by shape.name
    line_table = None
    summary_table = None
    found_names = []

    for shape in slide.shapes:
        if not shape.has_table:
            continue
        name = getattr(shape, "name", "")
        found_names.append(name)
        if name == "LineItems":
            line_table = shape.table
        elif name == "BillingSummary":
            summary_table = shape.table

    if line_table is None or summary_table is None:
        st.error(
            "❗ Could not find one or both tables by Name (shape.name).\n\n"
            "Tables found on Slide 1 have these names:\n\n"
            + "\n".join(f"- '{t}'" for t in found_names)
            + "\n\nPlease open your PPT template, right-click the correct table → Selection Pane, and set the table's name to exactly 'LineItems' or 'BillingSummary'."
        )
        st.stop()

    # 2) Determine Poppins 12 pt from the first blank row (row 1, col 0)
    max_rows = len(line_table.rows)
    if max_rows > 1:
        sample_cell = line_table.rows[1].cells[0]
    else:
        sample_cell = line_table.rows[0].cells[0]
    sample_para = sample_cell.text_frame.paragraphs[0]
    sample_run = sample_para.runs[0] if sample_para.runs else None

    if sample_run:
        base_font_name = sample_run.font.name or "Poppins"
        base_font_size = sample_run.font.size or Pt(12)
    else:
        base_font_name = "Poppins"
        base_font_size = Pt(12)

    # 3) Clear existing data rows (rows 1…end), leave header row 0 untouched
    available_data_rows = max_rows - 1
    for r_idx in range(1, max_rows):
        for c_idx in range(len(line_table.columns)):
            cell = line_table.rows[r_idx].cells[c_idx]
            cell.text = ""
            # add a blank run so Poppins 12 pt stays in place
            para = cell.text_frame.paragraphs[0]
            run = para.add_run()
            run.font.name = base_font_name
            run.font.size = base_font_size
            para.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    # 4) Fill each data row into rows 1… up to available_data_rows
    rows_to_fill = min(len(rows), available_data_rows)
    for i in range(rows_to_fill):
        row_data = rows[i]
        target_row = line_table.rows[i + 1]  # +1 to skip header
        for col_idx, key in enumerate(
            ["No.", "Item Description", "Weight", "Rate (₹)", "Amount (₹)"]
        ):
            txt = str(row_data.get(key, "")) if row_data.get(key) is not None else ""
            cell = target_row.cells[col_idx]
            cell.text = ""
            para = cell.text_frame.paragraphs[0]
            run = para.add_run()
            run.text = txt
            run.font.name = base_font_name
            run.font.size = base_font_size
            para.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    # 5) Compute Subtotal, Rounding, Net Payable
    amounts = []
    for row_data in rows:
        amt = pd.to_numeric(row_data.get("Amount (₹)", 0), errors="coerce")
        if pd.isna(amt):
            amt = 0.0
        amounts.append(float(amt))
    subtotal = sum(amounts)
    rounded_total = float(round(subtotal))
    rounding_value = rounded_total - subtotal
    net_payable = rounded_total

    # 6) Fill BillingSummary table (row 1→Subtotal, row 2→Rounding, row 3→NET PAYABLE)
    def set_summary_cell(r_idx, value):
        cell = summary_table.rows[r_idx].cells[1]
        cell.text = ""
        para = cell.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = f"{value:,.2f}"
        run.font.name = base_font_name
        run.font.size = base_font_size
        para.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

    set_summary_cell(1, subtotal)
    set_summary_cell(2, rounding_value)
    set_summary_cell(3, net_payable)

    # 7) Finally, save PPTX into a BytesIO and return the bytes
    out = BytesIO()
    prs.save(out)
    return out.getvalue()


# ──────────────────────────────────────────────────────────────────────────────
#  Streamlit App
# ──────────────────────────────────────────────────────────────────────────────
def main():
    st.title("💎 Rishab Gems ‒ Diamond Jewellery Invoice Generator")
    st.markdown(
        """
        Fill in each “Line Item” below. Tap **➕ Add Another Row** to append more items.
        Once you’re done, click **Generate Invoice** to download a fully-filled PowerPoint (.pptx).
        """
    )

    # Initialize “rows” in session state if not present
    if "rows" not in st.session_state:
        st.session_state.rows = [
            {"No.": "", "Item Description": "", "Weight": "", "Rate (₹)": "", "Amount (₹)": ""}
        ]

    rows = st.session_state.rows

    st.markdown("---")

    # For each row, present an expander "Line Item {i+1}" to group the five inputs vertically
    for idx in range(len(rows)):
        # Ensure each row has a unique key by assigning a UUID if not present
        if "uuid" not in rows[idx]:
            import uuid
            rows[idx]["uuid"] = str(uuid.uuid4())
        row_key = rows[idx]["uuid"]
        with st.expander(f"Line Item {idx+1}", expanded=True):
            no_val = st.text_input(
                label="No. (positive integer)",
                value=rows[idx]["No."],
                placeholder="1",
                key=f"No_{row_key}",
            )
            desc_val = st.text_input(
                label="Item Description (required)",
                value=rows[idx]["Item Description"],
                placeholder="Diamond Ring, Necklace…",
                key=f"Desc_{row_key}",
            )
            weight_val = st.text_input(
                label="Weight (gm) (non-negative)",
                value=rows[idx]["Weight"],
                placeholder="e.g. 1.25",
                key=f"Weight_{row_key}",
            )
            rate_val = st.text_input(
                label="Rate (₹) (non-negative)",
                value=rows[idx]["Rate (₹)"],
                placeholder="e.g. 45000",
                key=f"Rate_{row_key}",
            )
            amount_val = st.text_input(
                label="Amount (₹) (non-negative)",
                value=rows[idx]["Amount (₹)"],
                placeholder="e.g. 56250",
                key=f"Amount_{row_key}",
            )

            # Write back to session state
            st.session_state.rows[idx]["No."] = no_val
            st.session_state.rows[idx]["Item Description"] = desc_val
            st.session_state.rows[idx]["Weight"] = weight_val
            st.session_state.rows[idx]["Rate (₹)"] = rate_val
            st.session_state.rows[idx]["Amount (₹)"] = amount_val

    st.markdown("---")

    col_add, col_gen = st.columns([1, 1])
    with col_add:
        if st.button("➕ Add Another Row"):
            st.session_state.rows.append(
                {"No.": "", "Item Description": "", "Weight": "", "Rate (₹)": "", "Amount (₹)": ""}
            )
            st.rerun()

    with col_gen:
        generate_button = st.button("🖨️ Generate Invoice")

    # When “Generate Invoice” is clicked:
    if generate_button:
        # 1) Filter out rows where all five fields are blank
        filtered_rows = []
        for r in st.session_state.rows:
            if any(str(v).strip() != "" for v in r.values()):
                filtered_rows.append(r.copy())

        if not filtered_rows:
            st.error("No data entered. Please fill at least one line item before generating the invoice.")
            return

        # 2) Stricter validation and formatting
        errors = []
        validated = []
        for idx, r in enumerate(filtered_rows, start=1):
            row_errs = []
            # Validate No. (positive integer)
            try:
                no_int = int(float(r["No."]))
                if no_int <= 0:
                    row_errs.append("No. must be a positive integer.")
            except:
                row_errs.append("No. must be a positive integer.")

            # Validate Item Description (non-empty)
            if not str(r["Item Description"]).strip():
                row_errs.append("Item Description cannot be empty.")

            # Validate Weight (non-negative float)
            try:
                w = float(r["Weight"])
                if w < 0:
                    row_errs.append("Weight must be non-negative.")
            except:
                row_errs.append("Weight must be a number (e.g. 1.25).")

            # Validate Rate (non-negative float)
            try:
                rt = float(r["Rate (₹)"])
                if rt < 0:
                    row_errs.append("Rate (₹) must be non-negative.")
            except:
                row_errs.append("Rate (₹) must be a number (e.g. 45000).")

            # Validate Amount (non-negative float)
            try:
                amt = float(r["Amount (₹)"])
                if amt < 0:
                    row_errs.append("Amount (₹) must be non-negative.")
            except:
                row_errs.append("Amount (₹) must be a number (e.g. 56250).")

            if row_errs:
                errors.append(f"Line Item {idx}: " + "; ".join(row_errs))
            else:
                # If all fields are valid, re-format them as strings
                validated.append(
                    {
                        "No.": str(int(float(r["No."]))),
                        "Item Description": r["Item Description"],
                        "Weight": f"{float(r['Weight']):.2f}",
                        "Rate (₹)": f"{float(r['Rate (₹)']):.2f}",
                        "Amount (₹)": f"{float(r['Amount (₹)']):.2f}",
                    }
                )

        if errors:
            st.error("Please fix these errors before generating the invoice:")
            for e in errors:
                st.write(f"- {e}")
            return

        # 3) Sort validated rows by “No.”
        validated.sort(key=lambda x: int(x["No."]))

        # 4) Generate the PPTX
        try:
            pptx_bytes = generate_filled_invoice(validated, "invoice_template.pptx")
        except Exception as e:
            st.error(f"Unexpected error during PPT generation: {e}")
            return

        # 5) Offer download button
        if pptx_bytes:
            st.success("✅ Invoice generated successfully.")
            st.download_button(
                label="📥 Download Filled PPTX",
                data=pptx_bytes,
                file_name="RishabGems_Invoice_Filled.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )


if __name__ == "__main__":
    main()
    main()
