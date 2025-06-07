# app.py

import streamlit as st
st.set_page_config(
    page_title="Rishab Gems üíé",
    page_icon="diamond.ico",
    layout="centered",
)
page_element="""
<style>
[data-testid="stHeader"] {
  background-color: rgba(0,0,0,0);
}
#MainMenu {visibility: hidden;}
footer {visibility: hidden;}
header {visibility: hidden;}
div._link_gzau3_10 {
    display: none !important;
}
.st-emotion-cache-1cpx1b6 a {display: none;}
.st-emotion-cache-1vze3mj {display: none;}
</style>
"""
st.markdown(page_element, unsafe_allow_html=True)

hide_elements_style = """
<style>
/* Hide profile container */
div._profileContainer_gzau3_53 {
    display: none !important;
}

/* Hide 'Hosted with Streamlit' footer */
a._container_gzau3_1._viewerBadge_nim44_23 {
    display: none !important;
}
</style>
"""

st.markdown(hide_elements_style, unsafe_allow_html=True)

hide_streamlit_style = """
                <style>
                div[data-testid="stToolbar"] {
                visibility: hidden;
                height: 0%;
                position: fixed;
                }
                div[data-testid="stDecoration"] {
                visibility: hidden;
                height: 0%;
                position: fixed;
                }
                div[data-testid="stStatusWidget"] {
                visibility: hidden;
                height: 0%;
                position: fixed;
                }
                #MainMenu {
                visibility: hidden;
                height: 0%;
                }
                header {
                visibility: hidden;
                height: 0%;
                }
                footer {
                visibility: hidden;
                height: 0%;
                }
                </style>
                """
st.markdown(hide_streamlit_style, unsafe_allow_html=True)


import pandas as pd
from io import BytesIO
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import re
from datetime import datetime, timedelta
import time
import importlib.util

# app.py
import drive


# Import num2words from number-to-words.py
from number_to_words import num2words

with open("app/style.css") as css:
    st.markdown(f'<style>{css.read()}</style>', unsafe_allow_html=True)

# ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ
#  Helper: Fill PPTX template with invoice data (finding tables by alt_text)
# ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ
def generate_filled_invoice(rows, template_path, bill_info, payment_method, amount_in_words):
    prs = Presentation(template_path)
    slide = prs.slides[0]

    # Fill text boxes (titles + values)
    text_fields = {
        "Bill No": bill_info.get("Bill No", ""),
        "Bill Date": bill_info.get("Bill Date", ""),
        "Due Date": bill_info.get("Due Date", ""),
        "Biller Name": bill_info.get("Biller Name", ""),
        "Client Address": bill_info.get("Client Address", ""),
        "Client Phone Number": bill_info.get("Client Phone Number", ""),
        "Client Email": bill_info.get("Client Email", ""),
        "Client Bill To": bill_info.get("Client Bill To", ""),
    }
    field_titles = {
        "Bill No": "Bill No: ",
        "Bill Date": "Bill Date: ",
        "Due Date": "Due Date: ",
        "Biller Name": "Biller Name: ",
        "Client Bill To": "Bill To: ",
        "Client Address": "Address: ",
        "Client Phone Number": "Phone: ",
        "Client Email": "Email ID: ",
    }
    for shape in slide.shapes:
        if shape.has_text_frame:
            name = getattr(shape, "name", "")
            if name in text_fields and name in field_titles:
                value = str(text_fields[name])
                if name == "Client Address":
                    value = value[:65]
                shape.text_frame.clear()
                p = shape.text_frame.paragraphs[0]
                run_title = p.add_run()
                run_title.text = field_titles[name]
                run_title.font.bold = True
                run_title.font.name = "Poppins"
                run_title.font.size = Pt(12)
                run_value = p.add_run()
                run_value.text = value
                run_value.font.bold = False
                run_value.font.name = "Poppins"
                run_value.font.size = Pt(12)
                shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Handle Payment Method Checkboxes
    checkbox_names = {
        "**Cash**": "Cash Check",
        "**NEFT / IMPS**": "NEFT Check",
        "**UPI**": "UPI Check",
        "**Cheque**": "Cheque Check",
    }
    for shape in slide.shapes:
        if shape.has_text_frame:
            name = getattr(shape, "name", "")
            if name in checkbox_names.values():
                if name == checkbox_names.get(payment_method):
                    shape.text = "‚úî"
                    # Center align horizontally and vertically
                    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
                    # Optionally set font for tick
                    for run in shape.text_frame.paragraphs[0].runs:
                        run.font.name = "Poppins"
                        run.font.size = Pt(12)  # Adjust size as needed
                else:
                    shape.text = ""

    # Fill tables (LineItems and BillingSummary)
    line_table = None
    summary_table = None
    for shape in slide.shapes:
        if shape.has_table:
            if getattr(shape, "name", "") == "LineItems":
                line_table = shape.table
            elif getattr(shape, "name", "") == "BillingSummary":
                summary_table = shape.table

    if line_table is None or summary_table is None:
        raise Exception("Could not find LineItems or BillingSummary table in template.")

    # Get font style from first data cell
    max_rows = len(line_table.rows)
    if max_rows > 1:
        sample_cell = line_table.rows[1].cells[0]
    else:
        sample_cell = line_table.rows[0].cells[0]
    sample_para = sample_cell.text_frame.paragraphs[0]
    sample_run = sample_para.runs[0] if sample_para.runs else None
    if sample_run:
        base_font_name = sample_run.font.name or "Poppins"
        base_font_size = sample_run.font.size or Pt(12)
    else:
        base_font_name = "Poppins"
        base_font_size = Pt(12)

    # Clear existing data rows (rows 1‚Ä¶end), leave row 0 intact
    available_data_rows = max_rows - 1
    for r_idx in range(1, max_rows):
        for c_idx in range(len(line_table.columns)):
            cell = line_table.rows[r_idx].cells[c_idx]
            cell.text = ""
            para = cell.text_frame.paragraphs[0]
            run = para.add_run()
            run.font.name = base_font_name
            run.font.size = base_font_size
            para.alignment = PP_ALIGN.CENTER

    # Write each data row into rows 1‚Ä¶up to available_data_rows
    rows_to_fill = min(len(rows), available_data_rows)
    for i in range(rows_to_fill):
        row_data = rows[i]
        target_row = line_table.rows[i + 1]  # +1 to skip header
        for col_idx, key in enumerate(
            ["No.", "Item Description", "Weight", "Rate (‚Çπ)", "Amount (‚Çπ)"]
        ):
            value = row_data.get(key, "")
            if key == "Amount (‚Çπ)":
                try:
                    amount_val = float(value)
                    txt = f"{amount_val:,.2f}"
                except Exception:
                    txt = str(value)
            else:
                txt = str(value) if value is not None else ""
            cell = target_row.cells[col_idx]
            cell.text = ""
            para = cell.text_frame.paragraphs[0]
            run = para.add_run()
            run.text = txt
            run.font.name = base_font_name
            run.font.size = base_font_size
            para.alignment = PP_ALIGN.CENTER

    # Compute Subtotal, Rounding, Net Payable
    amounts = []
    for row_data in rows:
        amt = pd.to_numeric(row_data.get("Amount (‚Çπ)", 0), errors="coerce")
        if pd.isna(amt):
            amt = 0.0
        amounts.append(float(amt))
    subtotal = sum(amounts)
    rounded_total = float(round(subtotal))
    rounding_value = rounded_total - subtotal
    net_payable = rounded_total

    # Get amount in words
    amount_in_words = "Rupees " + num2words(net_payable) + " Only."

    # Fill BillingSummary table (row 1: Subtotal, row 2: Rounding, row 3: NET PAYABLE)
    def set_summary_cell(r_idx, value, prefix=""):
        cell = summary_table.rows[r_idx].cells[1]
        cell.text = ""
        para = cell.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = f"{prefix}{value:,.2f}"
        run.font.name = base_font_name
        run.font.size = base_font_size
        para.alignment = PP_ALIGN.CENTER

    set_summary_cell(1, subtotal)
    set_summary_cell(2, rounding_value)
    set_summary_cell(3, net_payable, prefix="‚Çπ ")

    # Write Amount In Words to the selection pane text box
    for shape in slide.shapes:
        if shape.has_text_frame and getattr(shape, "name", "") == "Amount In Words":
            shape.text_frame.clear()
            p = shape.text_frame.paragraphs[0]
            run = p.add_run()
            run.text = amount_in_words
            run.font.name = "Poppins"
            run.font.size = Pt(11)
            run.font.italic = True
            p.alignment = PP_ALIGN.CENTER  # Optional: center align text
            shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Save PPTX to memory
    out = BytesIO()
    prs.save(out)
    return out.getvalue()


# ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ
#  Streamlit App
# ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ‚îÄ
def parse_number(val):
    """
    Parse a number from a string, stripping spaces, commas, Rs, ‚Çπ, etc.
    Returns float or raises ValueError.
    """
    if val is None:
        raise ValueError("Empty value")
    s = str(val)
    s = s.replace(",", "")
    s = re.sub(r"(rs\.?|‚Çπ)", "", s, flags=re.IGNORECASE)
    s = s.strip()
    if not s:
        raise ValueError("Empty value")
    return float(s)

@st.cache_data
def slow_function():
    import time
    time.sleep(2)  # Simulating a slow operation
    return "Result from slow_function"

def main():
    # --- Minimalistic Login ---
    CORRECT_PIN = st.secrets.get("login_pin", "123456")
    if not st.session_state.get("authenticated"):
                # Centered logo above login
        st.markdown(
            """
            <div style="display: flex; flex-direction: column; align-items: center; margin-bottom: 1.5rem;">
                <img src="https://raw.githubusercontent.com/kingrishabdugar/RishabGems/744b262046e8c3152c21e4430da6ca59fd752e70/logo.png" width="80" style="border-radius:12px; box-shadow:0 2px 8px 0 rgba(0,0,0,0.06);" />
            </div>
            <div style="display: flex; flex-direction: row; align-items: center; justify-content: center; margin-bottom:0.5rem;">
                <img src="https://raw.githubusercontent.com/kingrishabdugar/RishabGems/refs/heads/main/image%20(1).png" width="50" style="margin-right: 0.5rem;" />
                <h1 style="margin: 0; padding: 0;">Login</h1>
            </div>
            """,
            unsafe_allow_html=True,
        )
        with open("app/style.css") as css:
            st.markdown(f'<style>{css.read()}</style>', unsafe_allow_html=True)
        pin = st.text_input(
            "PIN",  # Non-empty label for accessibility
            type="password",
            max_chars=6,
            key="pin",
            placeholder="_ _ _ _ _ _",
            label_visibility="collapsed",  # Hides it visually
            help="Enter your 6-digit PIN",
        )
        if st.button("![icon](https://raw.githubusercontent.com/kingrishabdugar/RishabGems/refs/heads/main/diamond.gif) **Login**", use_container_width=True):
            if pin == CORRECT_PIN:
                st.session_state["authenticated"] = True
                st.rerun()
            else:
                st.error("Incorrect PIN. Please try again.")
        return  # Stop here if not authenticated

    # --- Your existing app code below ---
    
    st.markdown(
        """
        <div style="display: flex; flex-direction: row; align-items: center; justify-content: center; margin-bottom:0.5rem;">
            <img src="https://raw.githubusercontent.com/kingrishabdugar/RishabGems/refs/heads/main/twindiamond.png" width="75" style="margin-right: 0.5rem;" />
            <h1 style="margin: 0; padding: 0;">Rishab Gems‚ÄíInvoice Generator</h1>
        </div>
        """,
        unsafe_allow_html=True,
    )

    st.markdown(
        """
        * Fill in your invoice line items below. Press **Add Another Row** to append more rows.
        * Once you‚Äôre done, click **Generate Invoice** to download a fully-filled PowerPoint (.pptx).
        * As soon as PPT is opened it can be shared by clicking on top right corner > 3 dots > Share as PDF > WhatsApp.
        """
    )

    # ‚îÄ‚îÄ NEW: Bill Info Section ‚îÄ‚îÄ
    today = datetime.today().date()
    default_due = today + timedelta(days=7)

    # Generate Bill No: e.g. RG-20240601-001
    if "bill_no" not in st.session_state:
        st.session_state["bill_no"] = f"RG-{today.strftime('%Y%m%d')}-{datetime.now().strftime('%H%M%S')}"

    # Use markdown for bold labels above each input
    st.markdown('<div class="subheading">Invoice Number</div>', unsafe_allow_html=True)
    bill_no = st.text_input(
        label="Invoice Number",
        key="bill_no",
        label_visibility="collapsed",
    )

    st.markdown('<div class="subheading">Bill Date</div>', unsafe_allow_html=True)
    bill_date = st.date_input("Bill Date", value=today, key="bill_date", label_visibility="collapsed")

    st.markdown('<div class="subheading">Due Date</div>', unsafe_allow_html=True)
    due_date = st.date_input("Due Date", value=default_due, key="due_date", label_visibility="collapsed")

    st.markdown('<div class="subheading">Biller Name</div>', unsafe_allow_html=True)
    biller_name = st.text_input("Biller Name", value="Mr. Manish Dugar", key="biller_name", label_visibility="collapsed")

    # st.markdown(f"**Invoice Number:** `{bill_no}`")

    # ‚îÄ‚îÄ NEW: Client Info Section ‚îÄ‚îÄ
    st.markdown('<div class="heading">Client Information</div>', unsafe_allow_html=True)

    st.markdown('<div class="subheading">Client Bill To</div>', unsafe_allow_html=True)
    client_bill_to = st.text_input(
        "Client Bill To",
        value="",
        placeholder="Enter client billing name",
        key="client_bill_to",
        label_visibility="collapsed",
    )

    st.markdown('<div class="subheading">Client Email</div>', unsafe_allow_html=True)
    client_email = st.text_input(
        "Client Email",
        value="",
        placeholder="Enter client email",
        key="client_email",
        label_visibility="collapsed",
    )
    
    st.markdown('<div class="subheading">Client Phone Number</div>', unsafe_allow_html=True)
    client_phone = st.text_input(
        "Client Phone Number",
        value="",
        placeholder="Enter client phone number",
        key="client_phone",
        label_visibility="collapsed",
    )

    st.markdown('<div class="subheading">Client Address (max 65 characters)</div>', unsafe_allow_html=True)
    client_address = st.text_input(
        "Client Address (max 65 characters)",
        value="",
        max_chars=65,
        placeholder="Enter client address",
        key="client_address",
        label_visibility="collapsed",
    )

    # ‚îÄ‚îÄ NEW: Payment Method Section ‚îÄ‚îÄ
    st.markdown('<div class="heading">Payment Method</div>', unsafe_allow_html=True)
    if "payment_method" not in st.session_state:
        st.session_state["payment_method"] = "**Cash**"
    
    st.markdown('<div class="subheading">Select Payment Method:</div>', unsafe_allow_html=True)
    payment_method = st.radio(
        "Select Payment Method:",
        options=["**Cash**", "**NEFT / IMPS**", "**UPI**", "**Cheque**"],
        index=0,
        key="payment_method",
        label_visibility = "collapsed"
    )

    # Initialize session-state rows
    if "rows" not in st.session_state:
        st.session_state.rows = [
            {"No.": "", "Item Description": "", "Weight": "", "Rate (‚Çπ)": "", "Amount (‚Çπ)": ""}
        ]

    rows = st.session_state.rows

    st.markdown("---")

    # Always use vertical card per line item, with label above each input (responsive & clear)
    for idx in range(len(rows)):
        with st.container():
            # Card-like background for each line item
            st.markdown(f'<div class="heading" style="margin-bottom:0.7rem;">Line Item {idx+1}</div>', unsafe_allow_html=True)

            # No.
            st.markdown('<div class="label-bold">No. (positive integer)</div>', unsafe_allow_html=True)
            default_no = rows[idx]["No."] if rows[idx]["No."] else str(idx + 1)
            no_val = st.text_input(
                label="No. (positive integer)",
                value=default_no,
                placeholder=str(idx + 1),
                key=f"No_{idx}",
                label_visibility="collapsed",
                disabled=True,  # Make it read-only
            )
            st.markdown("<div style='height: 0.3rem'></div>", unsafe_allow_html=True)

            # Item Description
            st.markdown('<div class="label-bold">Item Description (required)</div>', unsafe_allow_html=True)
            desc_val = st.text_input(
                label="Item Description (required)",
                value=rows[idx]["Item Description"],
                placeholder="Diamond Ring, Necklace‚Ä¶",
                key=f"Desc_{idx}",
                label_visibility="collapsed",
            )
            st.markdown("<div style='height: 0.3rem'></div>", unsafe_allow_html=True)

            # Per-row Weight Unit
            prev_unit = rows[idx-1]["Weight Unit"] if idx > 0 and "Weight Unit" in rows[idx-1] else "**carats**"
            st.markdown('<div class="label-bold">Weight Unit</div>', unsafe_allow_html=True)
            weight_unit = st.radio(
                f"Select Weight Unit for Row {idx+1}:",
                options=["**carats**", "**gms**"],
                index=0 if rows[idx].get("Weight Unit", prev_unit).replace("**", "") == "carats" else 1,
                key=f"weight_unit_{idx}",
                label_visibility="collapsed",
                horizontal=True,
            )
            # Save the selected unit for this row
            # Remove ** from weight unit for storage
            weight_unit = weight_unit.replace("**", "")
            rows[idx]["Weight Unit"] = weight_unit

            # Weight input (show unit)
            st.markdown(f'<div class="label-bold">Weight ({weight_unit}, non-negative)</div>', unsafe_allow_html=True)
            weight_val = st.text_input(
                label=f"Weight ({weight_unit}, non-negative)",
                value=rows[idx]["Weight"],
                placeholder=f"e.g. 1.25",
                key=f"Weight_{idx}",
                label_visibility="collapsed",
            )
            st.markdown("<div style='height: 0.3rem'></div>", unsafe_allow_html=True)

            # Rate (‚Çπ)
            st.markdown('<div class="label-bold">Rate (‚Çπ) (non-negative)</div>', unsafe_allow_html=True)
            rate_val = st.text_input(
                label="Rate (‚Çπ) (non-negative)",
                value=rows[idx]["Rate (‚Çπ)"],
                placeholder="e.g. 45000",
                key=f"Rate_{idx}",
                label_visibility="collapsed",
            )
            st.markdown("<div style='height: 0.3rem'></div>", unsafe_allow_html=True)  # Add spacing after Rate

            # --- Auto-calculate Amount ---
            try:
                w = float(weight_val)
                r = float(rate_val)
                auto_amount = f"{w * r:.2f}"
            except Exception:
                auto_amount = ""

            # Always auto-calculate and overwrite Amount (‚Çπ)
            amount_val = auto_amount

            # Amount (‚Çπ) (auto-calculated)
            st.markdown('<div class="label-bold">Amount (‚Çπ) (auto-calculated)</div>', unsafe_allow_html=True)
            st.text_input(
                label="Amount (‚Çπ) (auto-calculated)",
                value=amount_val,
                placeholder=auto_amount,
                key=f"Amount_{idx}",
                label_visibility="collapsed",
                disabled=False,  # Make it read-only
            )

            # Save back to session state
            st.session_state.rows[idx]["No."] = no_val
            st.session_state.rows[idx]["Item Description"] = desc_val
            st.session_state.rows[idx]["Weight"] = weight_val
            st.session_state.rows[idx]["Rate (‚Çπ)"] = rate_val
            st.session_state.rows[idx]["Amount (‚Çπ)"] = amount_val
            st.session_state.rows[idx]["Weight Unit"] = weight_unit

    st.markdown("---")

    # Centered, stacked buttons
    center_col = st.columns([3, 2, 3])[1]
    with center_col:
        if st.button("![icon](https://raw.githubusercontent.com/kingrishabdugar/RishabGems/refs/heads/main/pink-diamond.png) **Add Another Row**", use_container_width=True):
            st.session_state.rows.append(
                {
                    "No.": "",
                    "Item Description": "",
                    "Weight": "",
                    "Rate (‚Çπ)": "",
                    "Amount (‚Çπ)": "",
                }
            )
            st.rerun()
        st.markdown("<div style='height: 0.5rem'></div>", unsafe_allow_html=True)  # Small vertical gap
        generate_button = st.button("![icon](https://raw.githubusercontent.com/kingrishabdugar/RishabGems/refs/heads/main/pink-diamond.png) **Generate Invoice**", use_container_width=True)

    # When ‚ÄúGenerate Invoice‚Äù is clicked:
    if generate_button:
        # 1) Filter out fully blank rows
        filtered_rows = []
        for r in st.session_state.rows:
            if any(str(v).strip() != "" for v in r.values()):
                filtered_rows.append(r.copy())

        if not filtered_rows:
            st.error("No data entered. Please fill at least one line item before generating the invoice.")
            return

        # 2) Stricter validation
        errors = []
        validated = []
        for idx, r in enumerate(filtered_rows, start=1):
            row_errs = []
            # Validate No. (positive integer)
            try:
                no_int = int(parse_number(r["No."]))
                if no_int <= 0:
                    row_errs.append("No. must be a positive integer.")
            except Exception:
                row_errs.append("No. must be a positive integer.")

            # Validate Item Description (non-empty)
            if not str(r["Item Description"]).strip():
                row_errs.append("Item Description cannot be empty.")

            # Validate Weight (non-negative float)
            try:
                w = parse_number(r["Weight"])
                if w < 0:
                    row_errs.append("Weight must be non-negative.")
            except Exception:
                row_errs.append("Weight must be a number (e.g. 1.25).")

            # Validate Rate (non-negative float)
            try:
                rt = parse_number(r["Rate (‚Çπ)"])
                if rt < 0:
                    row_errs.append("Rate (‚Çπ) must be non-negative.")
            except Exception:
                row_errs.append("Rate (‚Çπ) must be a number (e.g. 45000).")

            # Validate Amount (non-negative float), allow blank and auto-calc
            amt_val = r["Amount (‚Çπ)"]
            amt_auto = False
            try:
                if str(amt_val).strip() == "":
                    # Try to auto-calculate
                    amt = parse_number(r["Weight"]) * parse_number(r["Rate (‚Çπ)"])
                    amt_auto = True
                else:
                    amt = parse_number(amt_val)
                if amt < 0:
                    row_errs.append("Amount (‚Çπ) must be non-negative.")
            except Exception:
                row_errs.append("Amount (‚Çπ) must be a number (e.g. 56250) or left blank for auto-calc.")

            if row_errs:
                errors.append(f"Row {idx}: " + "; ".join(row_errs))
            else:
                validated.append(
                    {
                        "No.": str(int(parse_number(r["No."]))),
                        "Item Description": r["Item Description"].strip(),
                        "Weight": f"{parse_number(r['Weight']):.2f} {r['Weight Unit']}",
                        "Rate (‚Çπ)": f"{parse_number(r['Rate (‚Çπ)']):.2f}",
                        "Amount (‚Çπ)": f"{amt:.2f}",
                    }
                )

        if errors:
            st.error("Please fix these errors before generating the invoice:")
            for e in errors:
                st.write(f"- {e}")
            return

        # 3) Sort by No.
        validated.sort(key=lambda x: int(x["No."]))

        # 4) Generate PPTX
        bill_info = {
            "Bill No": bill_no,
            "Bill Date": bill_date.strftime("%d-%m-%Y"),
            "Due Date": due_date.strftime("%d-%m-%Y"),
            "Biller Name": biller_name,
            "Client Address": client_address,
            "Client Phone Number": client_phone,
            "Client Email": client_email,
            "Client Bill To": client_bill_to,
        }

        # Calculate net_payable (same as in generate_filled_invoice)
        amounts = [float(x["Amount (‚Çπ)"]) for x in validated]
        subtotal = sum(amounts)
        net_payable = float(round(subtotal))

        # Get amount in words
        amount_in_words = "Rupees " + num2words(net_payable) + " Only."

        # Show Amount In Words textbox in UI (editable)
        st.markdown('<div class="subheading">Amount In Words</div>', unsafe_allow_html=True)
        amount_in_words = st.text_input(
            "Amount In Words",
            value=amount_in_words,
            key="amount_in_words",
            help="Amount in words for the invoice. You can edit if needed.",
            label_visibility="collapsed",
        )

        try:
            pptx_bytes = generate_filled_invoice(
                validated, "invoice_template.pptx", bill_info, payment_method, amount_in_words
            )
        except Exception as e:
            st.error(f"Unexpected error during PPT generation: {e}")
            return

        filename = f"Rishab_Gems_{client_bill_to}_{bill_no}_{client_phone}_{today.strftime('%Y%m%d')}.pptx"
        filename = re.sub(r'[\\/*?:"<>|]', "", filename)

        # Store in session state for use after rerun
        st.session_state["pptx_bytes"] = pptx_bytes
        st.session_state["pptx_filename"] = filename

        # ‚Üê NEW: Upload to Drive
        import traceback  # Add at the top of your file if not already imported

        print("Attempting to get folder ID for 'Tax Invoice'...")  # Debug
        folder_id = drive.get_folder_id("Tax Invoice")
        print(f"Folder ID: {folder_id}")  # Debug

        pptx_drive_fid = None
        pdf_drive_fid = None
        pdf_bytes = None
        pdf_filename = filename.replace(".pptx", ".pdf")

        if folder_id:
            try:
                # Show loading GIF before upload
                loading_gif_url = "https://raw.githubusercontent.com/kingrishabdugar/RishabGems/refs/heads/main/rotating-ring.webp"
                loading_placeholder = st.empty()
                loading_placeholder.markdown(
                    f'<div style="display:flex;justify-content:center;"><img src="{loading_gif_url}" alt="Loading..." style="height:400px;" /></div>',
                    unsafe_allow_html=True
                )
                print("Uploading PPTX to Drive...")  # Debug
                pptx_drive_fid = drive.upload_bytes_to_drive(
                    pptx_bytes,
                    filename,
                    folder_id,
                )
                print(f"Upload returned File ID: {pptx_drive_fid}")  # Debug
                st.success(f"‚úÖ PPTX uploaded to Drive (File ID: {pptx_drive_fid})")
                try:
                    # Show loading GIF before conversion
                    print("Converting PPTX to Google Slides...")  # Debug
                    slides_fid = drive.convert_pptx_to_slides(pptx_drive_fid, filename, folder_id)
                    print(f"Converted to Google Slides File ID: {slides_fid}")  # Debug

                    print("Attempting to export as PDF...")  # Debug
                    pdf_bytes = drive.export_drive_file_as_pdf(slides_fid)
                    print("PDF export successful")  # Debug
                    st.session_state["pdf_bytes"] = pdf_bytes
                    st.session_state["pdf_filename"] = pdf_filename
                    # Upload PDF to Drive
                    pdf_drive_fid = drive.upload_bytes_to_drive(
                        pdf_bytes,
                        pdf_filename,
                        folder_id,
                        mime_type="application/pdf"
                    )
                    st.success(f"‚úÖ PDF uploaded to Drive (File ID: {pdf_drive_fid})")
                except Exception as e:
                    loading_placeholder.empty()
                    st.error(f"Could not export as PDF: {e}")
                    st.write(traceback.format_exc())  # Print full traceback
            except Exception as e:
                loading_placeholder.empty()
                st.error(f"‚ö†Ô∏è Upload to Drive failed: {e}")
                st.write(traceback.format_exc())  # Print full traceback
        else:
            st.error("‚ö†Ô∏è 'Tax Invoice' folder not found in Drive. Share it with your service account and retry.")

    # Show the Share as PDF button if PDF is available (in place of Download PPTX)
    if "pdf_bytes" in st.session_state and "pdf_filename" in st.session_state:
        # Use custom HTML for better Android compatibility (auto download + share)
        # Streamlit's download_button does not support triggering the share sheet directly,
        # but setting download attribute helps with auto-download.
        b64_pdf = None
        try:
            import base64
            b64_pdf = base64.b64encode(st.session_state["pdf_bytes"]).decode()
        except Exception:
            pass

        if b64_pdf:
            # Use your preferred image URL here
            loading_placeholder.empty()  # Remove GIF 
            share_img_url = "https://raw.githubusercontent.com/kingrishabdugar/RishabGems/refs/heads/main/pink-diamond.png"
            custom_button = f"""
            <a href="data:application/pdf;base64,{b64_pdf}" download="{st.session_state['pdf_filename']}" class="golden-button" id="sharepdfbtn">
                <img src="{share_img_url}" alt="Share" style="height:1.5em;vertical-align:middle;margin-right:0.5em;" />
                Share as PDF
            </a>
            <script>
            // Try to auto-click for Android auto-download
            setTimeout(function() {{
                var btn = document.getElementById('sharepdfbtn');
                if(btn) btn.click();
            }}, 500);
            </script>
            """
            st.markdown(custom_button, unsafe_allow_html=True)
        else:
            st.download_button(
                label="**üì§ Share as PDF**",
                data=st.session_state["pdf_bytes"],
                file_name=st.session_state["pdf_filename"],
                mime="application/pdf",
            )        

if __name__ == "__main__":
    main()