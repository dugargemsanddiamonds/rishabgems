import streamlit as st
import pandas as pd
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment
from io import BytesIO
from matplotlib import pyplot as plt
from matplotlib import font_manager
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import tempfile
from fpdf import FPDF

# ─── CONSTANTS ────────────────────────────────────────────────────────────────
FONT_NAME = "Poppins"
FONT_SIZE_PT = 12
HEADER_FILL_COLOR = "#DDDDDD"    # header row color (must include '#')
ALT_ROW_FILL = "#F5F5F5"         # even rows light gray
ODD_ROW_FILL = "#FFFFFF"         # odd rows white
EXCEL_FILENAME = "filled_invoice.xlsx"
PNG_FILENAME   = "table_rendered.png"
PPTX_TEMPLATE  = "RISHAB-GEMS-Diamond-Jewellery-Invoice.pptx"

# The exact column headers in your template table:
COLUMNS = ["No.", "Item Description", "Weight", "Rate (₹)", "Amount (₹)"]

# ─── FUNCTIONS TO BUILD & STYLE EXCEL ────────────────────────────────────────
def build_and_style_excel(data_rows: list[list[str]]):
    # 1) Build DataFrame
    df = pd.DataFrame(data_rows, columns=COLUMNS)

    # 2) Write to BytesIO via ExcelWriter
    with BytesIO() as b_io:
        with pd.ExcelWriter(b_io, engine="openpyxl") as writer:
            df.to_excel(writer, index=False, sheet_name="Invoice")
            workbook = writer.book
        b_io.seek(0)
        wb = openpyxl.load_workbook(b_io)

    ws = wb["Invoice"]

    # 3) Style the header row
    header_font = Font(name=FONT_NAME, size=FONT_SIZE_PT, bold=True)
    header_fill = PatternFill("solid", fgColor=HEADER_FILL_COLOR.lstrip('#'))
    for col_idx, col_name in enumerate(COLUMNS, start=1):
        cell = ws.cell(row=1, column=col_idx)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center", vertical="center")

    # 4) Style data rows
    data_font = Font(name=FONT_NAME, size=FONT_SIZE_PT)
    for r in range(2, 2 + len(data_rows)):
        fill_color = ALT_ROW_FILL.lstrip('#') if (r % 2 == 0) else ODD_ROW_FILL.lstrip('#')
        for c in range(1, len(COLUMNS) + 1):
            cell = ws.cell(row=r, column=c)
            cell.font = data_font
            cell.fill = PatternFill("solid", fgColor=fill_color)
            if c == 2:  # "Item Description" column
                cell.alignment = Alignment(horizontal="left", vertical="center")
            else:
                cell.alignment = Alignment(horizontal="center", vertical="center")

    # 5) Autofit columns
    for col_idx, col_name in enumerate(COLUMNS, start=1):
        max_len = max(
            len(str(ws.cell(row=r, column=col_idx).value or ""))
            for r in range(1, 2 + len(data_rows))
        ) + 2
        ws.column_dimensions[openpyxl.utils.get_column_letter(col_idx)].width = max_len

    # 6) Adjust row heights
    for r in range(2, 2 + len(data_rows)):
        ws.row_dimensions[r].height = 20

    # 7) Save to disk
    wb.save(EXCEL_FILENAME)
    return wb

# ─── FUNCTION TO RENDER THE EXCEL SHEET AS A PNG ─────────────────────────────
def excel_to_png():
    df = pd.read_excel(EXCEL_FILENAME, sheet_name=0, dtype=str)

    # Create figure sized by rows & columns
    plt.close("all")
    fig, ax = plt.subplots(
        figsize=(len(df.columns) * 1.5, (len(df.index) + 1) * 0.5), dpi=200
    )
    ax.axis("off")

    # Build cell colors array
    n_rows, n_cols = df.shape
    cell_colours = []
    # Header row colors
    cell_colours.append([HEADER_FILL_COLOR] * n_cols)
    # Data rows alternating
    for r in range(n_rows):
        fill = ALT_ROW_FILL if ((r + 2) % 2 == 0) else ODD_ROW_FILL
        cell_colours.append([fill] * n_cols)

    # Create table in Matplotlib
    table = ax.table(
        cellText=[df.columns.tolist()] + df.values.tolist(),
        cellColours=cell_colours,
        cellLoc="center",
        colLabels=None,
        rowLabels=None,
        loc="center",
    )

    # Style each cell for fonts
    for (r, c), cell in table.get_celld().items():
        if r == 0:
            cell.get_text().set_fontweight("bold")
        cell.get_text().set_fontfamily(FONT_NAME)
        cell.get_text().set_fontsize(FONT_SIZE_PT - 2)
        if c == 1 and r > 0:
            cell.get_text().set_ha("left")
        else:
            cell.get_text().set_ha("center")
        cell.set_edgecolor("none")

    for spine in ax.spines.values():
        spine.set_visible(False)

    fig.tight_layout(pad=0)
    fig.savefig(PNG_FILENAME, bbox_inches="tight", transparent=False)
    plt.close(fig)

# ─── FUNCTION TO SUPERIMPOSE PNG ON EXISTING PPT TABLE LOCATION ──────────────
def insert_png_over_table(pptx_path: str, output_path: str, png_path: str):
    prs = Presentation(pptx_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            tbl = shape.table
            header_texts = [tbl.cell(0, c).text.strip() for c in range(len(COLUMNS))]
            if header_texts == COLUMNS:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                slide.shapes.add_picture(png_path, left, top, width=width, height=height)
                prs.save(output_path)
                return True
    return False

# ─── FUNCTION TO FILL "NET PAYABLE" IN SUMMARY ────────────────────────────────
def fill_net_payable(pptx_path: str, output_path: str, total_amount: float):
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_table:
                continue
            tbl = shape.table
            try:
                if len(tbl.rows) == 0:
                    continue
                last_row = tbl.rows[-1]
            except IndexError:
                continue  # Skip tables with no rows
            if any("NET PAYABLE" in cell.text.upper() for cell in last_row.cells):
                cell_to_fill = last_row.cells[1]
                cell_to_fill.text = f"₹ {total_amount:.2f}"
                for paragraph in cell_to_fill.text_frame.paragraphs:
                    paragraph.font.name = FONT_NAME
                    paragraph.font.size = Pt(FONT_SIZE_PT)
                    paragraph.alignment = PP_ALIGN.CENTER
                prs.save(output_path)
                return True
    return False

# ─── FUNCTION TO MAKE A SIMPLE PDF ────────────────────────────────────────────
def make_simple_pdf(pptx_path: str, output_pdf_path: str):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=14)
    pdf.cell(200, 10, txt="📄 Rishab Gems Invoice PDF", ln=True, align="C")
    pdf.ln(10)
    pdf.set_font("Arial", size=12)
    pdf.cell(200, 10, txt="(For a fully visual PDF of the slide, you’d need a PPT→PDF converter)", ln=True, align="L")
    pdf.output(output_pdf_path)

# ─── STREAMLIT UI ────────────────────────────────────────────────────────────
st.set_page_config(page_title="Rishab Gems Invoice Generator", layout="wide")
st.title("💎 Rishab Gems ‒ Diamond Jewellery Invoice Generator")
st.markdown(
    "Fill in the invoice rows below. When you click **Generate**, we will:\n\n"
    "1. Build a styled Excel with Poppins 12, bold headers, and alternate grey/white rows.  \n"
    "2. Convert that Excel into a high-res PNG.  \n"
    "3. Superimpose the PNG exactly onto your PPT template’s table.  \n"
    "4. Compute Net Payable and insert it into the summary table.  \n"
    "5. Let you download (a) the updated PPT and (b) a one-page reminder PDF."
)

if "rows" not in st.session_state:
    st.session_state.rows = [["", "", "", "", ""]]

for i in range(len(st.session_state.rows)):
    st.markdown(f"### Item Row {i+1}")
    cols = st.columns(5)
    for j in range(5):
        label = COLUMNS[j]
        current_val = st.session_state.rows[i][j]
        st.session_state.rows[i][j] = cols[j].text_input(f"{label} (Row {i+1})", value=current_val, key=f"r{i}c{j}")

if st.button("➕ Add Another Row"):
    st.session_state.rows.append(["", "", "", "", ""])

if st.button("🖨️ Generate Invoice"):
    try:
        # 1) Build & style Excel
        wb = build_and_style_excel(st.session_state.rows)

        # 2) Render Excel to PNG
        excel_to_png()

        # 3) Copy PPT template → temp, superimpose PNG, fill Net Payable
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as temp_ppt1:
            temp_ppt1.close()
            prs_template = Presentation(PPTX_TEMPLATE)
            prs_template.save(temp_ppt1.name)

            success = insert_png_over_table(
                pptx_path=temp_ppt1.name,
                output_path=temp_ppt1.name,
                png_path=PNG_FILENAME
            )
            if not success:
                st.error("Could not find a table in the PPTX whose header matches exactly " + str(COLUMNS))
                st.stop()

            total_amount = 0.0
            for row in st.session_state.rows:
                amt_str = row[4].strip().replace(",", "")
                if amt_str.replace(".", "", 1).isdigit():
                    total_amount += float(amt_str)
            success2 = fill_net_payable(
                pptx_path=temp_ppt1.name,
                output_path=temp_ppt1.name,
                total_amount=total_amount
            )
            if not success2:
                st.warning("Was unable to find a “NET PAYABLE” row in the summary. Check your template.")

        # 4) Offer final PPT
        with open(temp_ppt1.name, "rb") as f_ppt:
            st.download_button(
                "📥 Download Updated PPT",
                data=f_ppt,
                file_name="RishabGems_Invoice_Filled.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

        # 5) Create & offer a dummy PDF
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as temp_pdf:
            temp_pdf.close()
            make_simple_pdf(temp_ppt1.name, temp_pdf.name)
            with open(temp_pdf.name, "rb") as f_pdf:
                st.download_button(
                    "📥 Download Reminder PDF",
                    data=f_pdf,
                    file_name="RishabGems_Invoice_Reminder.pdf",
                    mime="application/pdf"
                )

        st.success("✅ All done! Your PPT now has a perfect Excel‐style table on top, and Net Payable is updated.")

    except Exception as e:
        st.error(f"⚠️ Encountered an error: {e}")
        raise
