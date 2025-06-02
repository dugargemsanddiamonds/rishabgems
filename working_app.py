# app.py

import streamlit as st
import pandas as pd
from io import BytesIO
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import re

# ──────────────────────────────────────────────────────────────────────────────
#  Helper: Fill PPTX template with invoice data (finding tables by alt_text)
# ──────────────────────────────────────────────────────────────────────────────
def generate_filled_invoice(rows, template_path):
    """
    1. Opens the PPTX template.
    2. Finds the "LineItems" table by Name (shape.name) on slide 0, writes data below header.
    3. Finds the "BillingSummary" table by Name on slide 0, fills Subtotal/Rounding/Net Payable.
    4. Returns pptx_bytes.
    """
    prs = Presentation(template_path)
    slide = prs.slides[0]

    # 1) Attempt to locate both tables by shape.name
    line_table = None
    summary_table = None
    found_names = []  # collect for debugging

    for shape in slide.shapes:
        if not shape.has_table:
            continue
        name = getattr(shape, "name", "")
        found_names.append(name)
        if name == "LineItems":
            line_table = shape.table
        elif name == "BillingSummary":
            summary_table = shape.table

    if line_table is None or summary_table is None:
        st.error(
            "❗ Could not find one or both tables by Name (shape.name).\n\n"
            "Tables found on Slide 1 have these names:\n\n"
            + "\n".join(f"- '{t}'" for t in found_names)
            + "\n\nPlease open your PPT template, right-click the correct table → Selection Pane, and set the table's name to exactly 'LineItems' or 'BillingSummary'."
        )
        st.stop()

    # 2) Grab Poppins 12 pt styling from row 1, col 0 (first blank row)
    max_rows = len(line_table.rows)
    if max_rows > 1:
        sample_cell = line_table.rows[1].cells[0]
    else:
        sample_cell = line_table.rows[0].cells[0]
    sample_para = sample_cell.text_frame.paragraphs[0]
    sample_run = sample_para.runs[0] if sample_para.runs else None
    if sample_run:
        base_font_name = sample_run.font.name or "Poppins"
        base_font_size = sample_run.font.size or Pt(12)
    else:
        base_font_name = "Poppins"
        base_font_size = Pt(12)

    # 3) Clear existing data rows (rows 1…end), leave row 0 intact
    available_data_rows = max_rows - 1
    for r_idx in range(1, max_rows):
        for c_idx in range(len(line_table.columns)):
            cell = line_table.rows[r_idx].cells[c_idx]
            cell.text = ""
            para = cell.text_frame.paragraphs[0]
            run = para.add_run()
            run.font.name = base_font_name
            run.font.size = base_font_size
            para.alignment = PP_ALIGN.CENTER  # Center align

    # 4) Write each data row into rows 1…up to available_data_rows
    rows_to_fill = min(len(rows), available_data_rows)
    for i in range(rows_to_fill):
        row_data = rows[i]
        target_row = line_table.rows[i + 1]  # +1 to skip header
        for col_idx, key in enumerate(
            ["No.", "Item Description", "Weight", "Rate (₹)", "Amount (₹)"]
        ):
            txt = str(row_data.get(key, "")) if row_data.get(key) is not None else ""
            cell = target_row.cells[col_idx]
            cell.text = ""
            para = cell.text_frame.paragraphs[0]
            run = para.add_run()
            run.text = txt
            run.font.name = base_font_name
            run.font.size = base_font_size
            para.alignment = PP_ALIGN.CENTER  # Center align

    # 5) Compute Subtotal, Rounding, Net Payable
    amounts = []
    for row_data in rows:
        amt = pd.to_numeric(row_data.get("Amount (₹)", 0), errors="coerce")
        if pd.isna(amt):
            amt = 0.0
        amounts.append(float(amt))
    subtotal = sum(amounts)
    rounded_total = float(round(subtotal))
    rounding_value = rounded_total - subtotal
    net_payable = rounded_total

    # 6) Fill BillingSummary table (row 1: Subtotal, row 2: Rounding, row 3: NET PAYABLE)
    def set_summary_cell(r_idx, value):
        cell = summary_table.rows[r_idx].cells[1]
        cell.text = ""
        para = cell.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = f"{value:,.2f}"
        run.font.name = base_font_name
        run.font.size = base_font_size
        para.alignment = PP_ALIGN.CENTER  # Center align

    set_summary_cell(1, subtotal)
    set_summary_cell(2, rounding_value)
    set_summary_cell(3, net_payable)

    # 7) Save PPTX to memory
    out = BytesIO()
    prs.save(out)
    return out.getvalue()


# ──────────────────────────────────────────────────────────────────────────────
#  Streamlit App
# ──────────────────────────────────────────────────────────────────────────────
def parse_number(val):
    """
    Parse a number from a string, stripping spaces, commas, Rs, ₹, etc.
    Returns float or raises ValueError.
    """
    if val is None:
        raise ValueError("Empty value")
    s = str(val)
    s = s.replace(",", "")
    s = re.sub(r"(rs\.?|₹)", "", s, flags=re.IGNORECASE)
    s = s.strip()
    if not s:
        raise ValueError("Empty value")
    return float(s)

def main():
    st.set_page_config(
        page_title="💎 Rishab Gems ‒ Diamond Jewellery Invoice Generator",
        layout="wide",
    )

    st.title("💎 Rishab Gems ‒ Diamond Jewellery Invoice Generator")
    st.markdown(
        """
        Fill in your invoice line items below. Press “➕ Add Another Row” to append more rows.
        Once you’re done, click **Generate Invoice** to download a fully-filled PowerPoint (.pptx).
        """
    )

    # Initialize session-state rows
    if "rows" not in st.session_state:
        st.session_state.rows = [
            {"No.": "", "Item Description": "", "Weight": "", "Rate (₹)": "", "Amount (₹)": ""}
        ]

    rows = st.session_state.rows

    st.markdown("---")

    # Responsive UI: Use vertical cards on mobile, horizontal rows on desktop
    is_mobile = st.session_state.get("_is_mobile", None)
    if is_mobile is None:
        # Simple user-agent check (Streamlit doesn't expose directly, so use a hack)
        is_mobile = st.query_params().get("mobile", ["0"])[0] == "1"
        st.session_state["_is_mobile"] = is_mobile

    if not is_mobile:
        # Desktop: horizontal row layout
        # Header labels (not inputs)
        header_cols = st.columns([1, 3, 2, 2, 2])
        header_cols[0].markdown("**No.**")
        header_cols[1].markdown("**Item Description**")
        header_cols[2].markdown("**Weight (gm)**")
        header_cols[3].markdown("**Rate (₹)**")
        header_cols[4].markdown("**Amount (₹)**")

        for idx in range(len(rows)):
            c1, c2, c3, c4, c5 = st.columns([1, 3, 2, 2, 2], gap="small")
            no_val = c1.text_input(
                label="No.",
                value=rows[idx]["No."],
                placeholder="1",
                key=f"No_{idx}",
                label_visibility="collapsed",
            )
            desc_val = c2.text_input(
                label="Item Description",
                value=rows[idx]["Item Description"],
                placeholder="Diamond Ring, Necklace…",
                key=f"Desc_{idx}",
                label_visibility="collapsed",
            )
            weight_val = c3.text_input(
                label="Weight (gm)",
                value=rows[idx]["Weight"],
                placeholder="e.g. 1.25",
                key=f"Weight_{idx}",
                label_visibility="collapsed",
            )
            rate_val = c4.text_input(
                label="Rate (₹)",
                value=rows[idx]["Rate (₹)"],
                placeholder="e.g. 45000",
                key=f"Rate_{idx}",
                label_visibility="collapsed",
            )
            amount_val = c5.text_input(
                label="Amount (₹)",
                value=rows[idx]["Amount (₹)"],
                placeholder="e.g. 56250",
                key=f"Amount_{idx}",
                label_visibility="collapsed",
            )
            st.session_state.rows[idx]["No."] = no_val
            st.session_state.rows[idx]["Item Description"] = desc_val
            st.session_state.rows[idx]["Weight"] = weight_val
            st.session_state.rows[idx]["Rate (₹)"] = rate_val
            st.session_state.rows[idx]["Amount (₹)"] = amount_val
    else:
        # Mobile: vertical card per line item
        for idx in range(len(rows)):
            with st.container():
                st.markdown(f"**Line Item {idx+1}**")
                no_val = st.text_input(
                    label="No. (positive integer)",
                    value=rows[idx]["No."],
                    placeholder="1",
                    key=f"No_{idx}",
                )
                desc_val = st.text_input(
                    label="Item Description (required)",
                    value=rows[idx]["Item Description"],
                    placeholder="Diamond Ring, Necklace…",
                    key=f"Desc_{idx}",
                )
                weight_val = st.text_input(
                    label="Weight (gm) (non-negative)",
                    value=rows[idx]["Weight"],
                    placeholder="e.g. 1.25",
                    key=f"Weight_{idx}",
                )
                rate_val = st.text_input(
                    label="Rate (₹) (non-negative)",
                    value=rows[idx]["Rate (₹)"],
                    placeholder="e.g. 45000",
                    key=f"Rate_{idx}",
                )
                amount_val = st.text_input(
                    label="Amount (₹) (non-negative, auto if blank)",
                    value=rows[idx]["Amount (₹)"],
                    placeholder="e.g. 56250",
                    key=f"Amount_{idx}",
                )
                st.session_state.rows[idx]["No."] = no_val
                st.session_state.rows[idx]["Item Description"] = desc_val
                st.session_state.rows[idx]["Weight"] = weight_val
                st.session_state.rows[idx]["Rate (₹)"] = rate_val
                st.session_state.rows[idx]["Amount (₹)"] = amount_val

    st.markdown("---")

    col_a, col_b = st.columns([1, 1])
    with col_a:
        if st.button("➕ Add Another Row"):
            st.session_state.rows.append(
                {
                    "No.": "",
                    "Item Description": "",
                    "Weight": "",
                    "Rate (₹)": "",
                    "Amount (₹)": "",
                }
            )
            st.rerun()

    with col_b:
        generate_button = st.button("🖨️ Generate Invoice")

    # When “Generate Invoice” is clicked:
    if generate_button:
        # 1) Filter out fully blank rows
        filtered_rows = []
        for r in st.session_state.rows:
            if any(str(v).strip() != "" for v in r.values()):
                filtered_rows.append(r.copy())

        if not filtered_rows:
            st.error("No data entered. Please fill at least one line item before generating the invoice.")
            return

        # 2) Stricter validation
        errors = []
        validated = []
        for idx, r in enumerate(filtered_rows, start=1):
            row_errs = []
            # Validate No. (positive integer)
            try:
                no_int = int(parse_number(r["No."]))
                if no_int <= 0:
                    row_errs.append("No. must be a positive integer.")
            except Exception:
                row_errs.append("No. must be a positive integer.")

            # Validate Item Description (non-empty)
            if not str(r["Item Description"]).strip():
                row_errs.append("Item Description cannot be empty.")

            # Validate Weight (non-negative float)
            try:
                w = parse_number(r["Weight"])
                if w < 0:
                    row_errs.append("Weight must be non-negative.")
            except Exception:
                row_errs.append("Weight must be a number (e.g. 1.25).")

            # Validate Rate (non-negative float)
            try:
                rt = parse_number(r["Rate (₹)"])
                if rt < 0:
                    row_errs.append("Rate (₹) must be non-negative.")
            except Exception:
                row_errs.append("Rate (₹) must be a number (e.g. 45000).")

            # Validate Amount (non-negative float), allow blank and auto-calc
            amt_val = r["Amount (₹)"]
            amt_auto = False
            try:
                if str(amt_val).strip() == "":
                    # Try to auto-calculate
                    amt = parse_number(r["Weight"]) * parse_number(r["Rate (₹)"])
                    amt_auto = True
                else:
                    amt = parse_number(amt_val)
                if amt < 0:
                    row_errs.append("Amount (₹) must be non-negative.")
            except Exception:
                row_errs.append("Amount (₹) must be a number (e.g. 56250) or left blank for auto-calc.")

            if row_errs:
                errors.append(f"Row {idx}: " + "; ".join(row_errs))
            else:
                validated.append(
                    {
                        "No.": str(int(parse_number(r["No."]))),
                        "Item Description": r["Item Description"].strip(),
                        "Weight": f"{parse_number(r['Weight']):.2f}",
                        "Rate (₹)": f"{parse_number(r['Rate (₹)']):.2f}",
                        "Amount (₹)": f"{amt:.2f}",
                    }
                )

        if errors:
            st.error("Please fix these errors before generating the invoice:")
            for e in errors:
                st.write(f"- {e}")
            return

        # 3) Sort by No.
        validated.sort(key=lambda x: int(x["No."]))

        # 4) Generate PPTX
        try:
            pptx_bytes = generate_filled_invoice(validated, "invoice_template.pptx")
        except Exception as e:
            st.error(f"Unexpected error during PPT generation: {e}")
            return

        # 5) Download button for PPTX
        if pptx_bytes:
            st.success("✅ Invoice generated successfully.")
            st.download_button(
                label="📥 Download Filled PPTX",
                data=pptx_bytes,
                file_name="RishabGems_Invoice_Filled.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )


if __name__ == "__main__":
    main()
